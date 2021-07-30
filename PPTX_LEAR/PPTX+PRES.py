from pptx import Presentation
import os


prs = Presentation()

layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(layout)
title = slide.shapes.title
title.text = 'This is my title'
subtitle = slide.placeholders[1]
subtitle.text = 'This is my subtitle'

slides = [

    ['Title of presentation',
     'Subtitle of presentation',
     0],
    ['Title of paragraph',
    'Content (bullet',
    1]


]

prs.save('MyPresentation.pptx')
os.startfile('MyPresentation.pptx')